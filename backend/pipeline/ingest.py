import os
import io
from pathlib import Path

import pdfplumber
from pptx import Presentation

from config import get_settings

settings = get_settings()
DOCUMENTS_PATH = settings.upload_dir


def read_txt_from_path(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def read_pdf_from_path(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    return text


def read_pptx_from_path(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# ── Works from bytes (used later in AWS Lambda) ──────────────────────────────

def read_pdf_from_bytes(file_bytes):
    text = ""
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    return text


def read_pptx_from_bytes(file_bytes):
    text = ""
    presentation = Presentation(io.BytesIO(file_bytes))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# ── Main extract functions ────────────────────────────────────────────────────

def extract_text(file_path):
    """Extract text from a local file path. Used for local development."""
    suffix = Path(file_path).suffix.lower()

    if suffix == ".txt":
        return read_txt_from_path(file_path)
    elif suffix == ".pdf":
        return read_pdf_from_path(file_path)
    elif suffix == ".pptx":
        return read_pptx_from_path(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def extract_text_from_bytes(file_bytes, extension):
    """Extract text from raw bytes. Used in AWS Lambda later."""
    extension = extension.lower()

    if extension == ".txt":
        return file_bytes.decode("utf-8")
    elif extension == ".pdf":
        return read_pdf_from_bytes(file_bytes)
    elif extension == ".pptx":
        return read_pptx_from_bytes(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {extension}")


def load_all_documents():
    """Load all documents from local documents/ folder."""
    documents = []

    for file in os.listdir(DOCUMENTS_PATH):
        path = os.path.join(DOCUMENTS_PATH, file)

        try:
            text = extract_text(path)
            documents.append({
                "file_name": file,
                "text": text
            })
            print(f"  Loaded: {file}")

        except Exception as e:
            print(f"  Skipping {file}: {e}")

    return documents
